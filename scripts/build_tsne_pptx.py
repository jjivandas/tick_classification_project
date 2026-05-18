"""Render t-SNE PNGs from coords.csv (replicating viewer.html styling) and bundle into a PowerPoint.

Primary label is always 'species'. Secondary labels cycle through the remaining metadata fields.
"""
from __future__ import annotations

import colorsys
import json
from pathlib import Path

import pandas as pd
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt

RUN_DIR = Path("results/tsne/runs/260504_234542")
OUT_DIR = Path("results/tsne/powerpoint_exports")
PNG_DIR = OUT_DIR / "pngs"
PNG_DIR.mkdir(parents=True, exist_ok=True)

SHOTS = [0, 1, 10, 25, 35]
METHOD = "tsne"
PRIMARY = "species"
SECONDARIES = ["species", "sex", "life_stage", "attached", "pathogen", "pathogen_result", "tick_condition"]

IMG_WIDTH = 1400
IMG_HEIGHT = 900

manifest = json.loads((RUN_DIR / "manifest.json").read_text())
species_palette: dict[str, str] = dict(manifest["species_palette"])
species_palette["Ixodes scapularis"] = "#4a148c"  # dark purple override (was #9467bd)
df_all = pd.read_csv(RUN_DIR / "coords.csv")


def safe_value(v) -> str:
    if pd.isna(v):
        return "(missing)"
    s = str(v).strip()
    return s if s else "(missing)"


def hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def clamp(v, lo, hi):
    return max(lo, min(hi, v))


def style_color(base_hex: str, sat_mul: float, light_delta: float, alpha: float) -> str:
    r, g, b = hex_to_rgb(base_hex)
    h, l, s = colorsys.rgb_to_hls(r / 255, g / 255, b / 255)
    s_pct = clamp(s * 100 * sat_mul, 8, 100)
    l_pct = clamp(l * 100 + light_delta, 10, 90)
    a = clamp(alpha, 0.28, 1.0)
    return f"hsla({h*360:.1f}, {s_pct:.1f}%, {l_pct:.1f}%, {a:.3f})"


def secondary_style(index: int, total: int) -> tuple[float, float, float]:
    denom = max(1, total - 1)
    t = index / denom
    return (1.4 - 0.7 * t, -22 + 44 * t, 0.98 - 0.18 * t)


def shot_label(k: int) -> str:
    return "0 (raw BioCLIP)" if k == 0 else str(k)


def make_figure(shots: int, secondary: str) -> go.Figure:
    sub = df_all[(df_all["shots"] == shots) & (df_all["method"] == METHOD)].copy()
    sub[PRIMARY] = sub[PRIMARY].map(safe_value)
    sub[secondary] = sub[secondary].map(safe_value)

    primary_vals = sorted(sub[PRIMARY].unique())
    secondary_vals = sorted(sub[secondary].unique())
    sec_index = {v: i for i, v in enumerate(secondary_vals)}

    fig = go.Figure()
    same_label = secondary == PRIMARY
    # group order: primary then secondary, alphabetical (matches JS)
    grouped = sub.groupby([PRIMARY, secondary], sort=True)
    for (p, s), g in grouped:
        base = species_palette.get(p, "#7f7f7f")
        if same_label:
            # species-only view: use base palette color directly, no tint modulation
            color = base
            name = p
        else:
            sat_mul, light_delta, alpha = secondary_style(sec_index[s], len(secondary_vals))
            color = style_color(base, sat_mul, light_delta, alpha)
            name = f"{p} | {s}"
        fig.add_trace(go.Scattergl(
            x=g["x"], y=g["y"], mode="markers",
            name=name,
            marker=dict(size=8, color=color, line=dict(width=0)),
            hoverinfo="skip",
        ))

    input_desc = "raw BioCLIP (512-D)" if shots == 0 else f"mean SVM predict_proba (5-D, {shots}-shot, 100 MC runs)"
    fig.update_layout(
        title=dict(
            text=f"t-SNE of {input_desc}<br><sup>Primary: {PRIMARY}  ·  Secondary: {secondary}  ·  N={len(sub)} specimens</sup>",
            font=dict(size=16),
            x=0.5, xanchor="center",
        ),
        legend=dict(title=dict(text="Species" if secondary == PRIMARY else "Primary | Secondary"), orientation="v"),
        xaxis=dict(title="t-SNE 1", zeroline=False, gridcolor="rgba(148,163,184,0.18)"),
        yaxis=dict(title="t-SNE 2", zeroline=False, gridcolor="rgba(148,163,184,0.18)"),
        paper_bgcolor="white", plot_bgcolor="white",
        width=IMG_WIDTH, height=IMG_HEIGHT,
        margin=dict(l=60, r=24, t=80, b=56),
    )
    return fig


def png_path(shots: int, secondary: str) -> Path:
    return PNG_DIR / f"tsne_k{shots}_species_x_{secondary}.png"


def render_all_pngs() -> None:
    for k in SHOTS:
        for sec in SECONDARIES:
            out = png_path(k, sec)
            fig = make_figure(k, sec)
            fig.write_image(out, format="png", width=IMG_WIDTH, height=IMG_HEIGHT, scale=2)
            print(f"wrote {out}")


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # fully blank layout

    def add_text_slide(text: str, font_size: int = 44) -> None:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), prs.slide_width - Inches(1.0), Inches(1.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = 2  # center
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = True

    def add_image_slide(heading: str, image: Path) -> None:
        slide = prs.slides.add_slide(blank)
        # heading
        hbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), prs.slide_width - Inches(0.6), Inches(0.6))
        tf = hbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = heading
        p.alignment = 2
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.bold = True
        # image, scaled to fit
        max_w = prs.slide_width - Inches(0.6)
        max_h = prs.slide_height - Inches(1.0)
        # plotly png is IMG_WIDTH x IMG_HEIGHT (scale=2 -> 2x pixels but same aspect)
        aspect = IMG_WIDTH / IMG_HEIGHT
        if max_w / max_h > aspect:
            h = max_h
            w = int(h * aspect)
        else:
            w = max_w
            h = int(w / aspect)
        left = (prs.slide_width - w) // 2
        top = Inches(0.85)
        slide.shapes.add_picture(str(image), left, top, width=w, height=h)

    # 1. Title slide
    add_text_slide("tsne plots tickid", font_size=54)

    # 2. Per-shot section + plots
    for k in SHOTS:
        add_text_slide(f"{k} shot" if k != 0 else "0 shot", font_size=54)
        for sec in SECONDARIES:
            heading = f"t-SNE ({shot_label(k)}-shot)  ·  primary: species  ·  secondary: {sec}"
            add_image_slide(heading, png_path(k, sec))

    out = OUT_DIR / "tsne_plots_tickid.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    render_all_pngs()
    out = build_pptx()
    print(f"\nDeck written: {out}")
